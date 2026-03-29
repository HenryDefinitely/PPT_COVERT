import streamlit as st
import cv2
import easyocr
import numpy as np
import io
import gc  # 新增：用於手動強制釋放記憶體
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

# ---------------------------------------------------------
# 1. 設定網頁標題與外觀
# ---------------------------------------------------------
st.set_page_config(page_title="圖片轉可編輯 PPT 神器", page_icon="🪄")
st.title("🪄 圖片轉可編輯 PPT 神器")
st.write("上傳一張圖片，AI 會自動將裡面的圖案轉成「積木色塊」，並把文字轉成「PPT 文字方塊」！")

# ---------------------------------------------------------
# 2. 載入 AI 模型 (只載入一次，避免重複撐爆 RAM)
# ---------------------------------------------------------
@st.cache_resource
def load_ocr_model():
    # 強制使用 CPU 模式以適應免費雲端伺服器
    return easyocr.Reader(['ch_tra', 'en'], gpu=False)

reader = load_ocr_model()

# ---------------------------------------------------------
# 2.5 新增：快取圖片辨識結果 (@st.cache_data)
# ---------------------------------------------------------
@st.cache_data(show_spinner=False)
def perform_ocr(image_bytes):
    """
    這個函數會把圖片的 Byte 資料與 OCR 結果綁定。
    只要上傳的圖片是同一張，就不會重複執行吃力的 PyTorch 運算！
    """
    img_np = cv2.imdecode(np.frombuffer(image_bytes, np.uint8), 1)
    results = reader.readtext(img_np)
    
    # 辨識完畢立刻釋放暫存矩陣
    del img_np
    gc.collect()
    return results

# ---------------------------------------------------------
# 3. 網頁控制面板 (側邊欄)
# ---------------------------------------------------------
with st.sidebar:
    st.header("⚙️ 轉換設定")
    resolution = st.selectbox(
        "選擇圖案精細度 (解析度)",
        ('Low (快, 150區塊)', 'Medium (適中, 250區塊)', 'High (慢, 350區塊)')
    )
    # 抓取對應的數值
    if 'Low' in resolution: res_val = 150
    elif 'Medium' in resolution: res_val = 250
    else: res_val = 350

# ---------------------------------------------------------
# 4. 圖片上傳區塊
# ---------------------------------------------------------
uploaded_file = st.file_uploader("📂 請上傳圖片 (支援 PNG, JPG, JPEG)", type=['png', 'jpg', 'jpeg'])

if uploaded_file is not None:
    st.image(uploaded_file, caption="上傳的圖片預覽", use_container_width=True)
    
    if st.button("🚀 開始轉換為 PPT"):
        progress_bar = st.progress(0)
        status_text = st.empty()
        
        # 取得圖片的二進位資料 (供快取函數比對使用)
        file_bytes = uploaded_file.getvalue()
        
        status_text.text("🧠 正在讀取圖片與辨識文字 (若為相同圖片將啟用秒速快取)...")
        
        # 呼叫快取的 OCR 函數
        text_results = perform_ocr(file_bytes)

        # 重新將檔案轉換為 OpenCV 格式供後續繪圖使用
        img = cv2.imdecode(np.frombuffer(file_bytes, np.uint8), 1)
        img_h, img_w = img.shape[:2]

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # === 【文字虛擬橡皮擦】 ===
        text_data_to_draw = []
        if text_results:
            for (bbox, text, prob) in text_results:
                text_data_to_draw.append((bbox, text))
                pts = np.array(bbox, np.int32)
                cv2.fillPoly(img, [pts], (255, 255, 255))

        # === 【繪製圖形 (掃描線積木填色)】 ===
        status_text.text("🎨 正在繪製積木色塊...")
        ratio = res_val / img_w
        block_h = int(img_h * ratio)
        small_img = cv2.resize(img, (res_val, block_h), interpolation=cv2.INTER_AREA)

        max_ppt_w, max_ppt_h = 720, 540
        ppt_scale = min(max_ppt_w / res_val, max_ppt_h / block_h) * 0.8
        offset_x = (max_ppt_w - (res_val * ppt_scale)) / 2
        offset_y = (max_ppt_h - (block_h * ppt_scale)) / 2

        for y in range(block_h):
            progress_bar.progress((y + 1) / block_h)
            x = 0
            while x < res_val:
                b, g, r = small_img[y, x]
                if r > 240 and g > 240 and b > 240:
                    x += 1
                    continue
                    
                start_x = x
                while x < res_val:
                    nb, ng, nr = small_img[y, x]
                    if abs(int(nr)-int(r)) <= 3 and abs(int(ng)-int(g)) <= 3 and abs(int(nb)-int(b)) <= 3:
                        x += 1
                    else:
                        break
                width = x - start_x
                
                ppt_x = start_x * ppt_scale + offset_x
                ppt_y = y * ppt_scale + offset_y
                ppt_width = width * ppt_scale
                ppt_height = ppt_scale
                
                rect = slide.shapes.add_shape(1, Pt(ppt_x), Pt(ppt_y), Pt(ppt_width + 0.3), Pt(ppt_height + 0.3))
                rect.fill.solid()
                rect.fill.fore_color.rgb = RGBColor(int(r), int(g), int(b))
                rect.line.fill.background() 

        # === 【在圖形上疊加文字方塊】 ===
        status_text.text("📝 正在排版文字方塊...")
        for bbox, text in text_data_to_draw:
            orig_x, orig_y = bbox[0][0], bbox[0][1]
            orig_w, orig_h = bbox[1][0] - bbox[0][0], bbox[2][1] - bbox[1][1]
            
            ppt_text_x = (orig_x / img_w) * (res_val * ppt_scale) + offset_x
            ppt_text_y = (orig_y / img_h) * (block_h * ppt_scale) + offset_y
            ppt_text_w = max((orig_w / img_w) * (res_val * ppt_scale), 10)
            ppt_text_h = max((orig_h / img_h) * (block_h * ppt_scale), 10)
            
            txBox = slide.shapes.add_textbox(Pt(ppt_text_x), Pt(ppt_text_y), Pt(ppt_text_w), Pt(ppt_text_h))
            tf = txBox.text_frame
            tf.text = text
            
            font_size = max(int(ppt_text_h * 0.7), 10)
            tf.paragraphs[0].font.size = Pt(font_size)
            tf.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)

        # === 【將 PPT 存入記憶體並提供下載】 ===
        status_text.text("✅ 轉換完成！請點擊下方按鈕下載 PPT。")
        
        ppt_stream = io.BytesIO()
        prs.save(ppt_stream)
        ppt_stream.seek(0)
        
        # 🛡️ 終極防護：強制清空龐大的圖形矩陣，歸還系統記憶體
        del img
        del small_img
        gc.collect()
        
        st.success("轉換成功！")
        st.download_button(
            label="📥 下載可編輯的 PPT 檔案",
            data=ppt_stream,
            file_name="converted_editable.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
